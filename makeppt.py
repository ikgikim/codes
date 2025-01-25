import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ppt_from_text_files():
    prs = Presentation()
    
    folder_path = os.getcwd()  # 현재 작업 디렉토리
    titles = []  # 제목과 페이지 번호를 저장할 리스트
    file_titles = []  # 파일 이름 저장

    for filename in os.listdir(folder_path):
        if filename.endswith('.txt'):
            file_titles.append(filename[:-4])  # 확장자 제거한 파일 이름
            file_path = os.path.join(folder_path, filename)
            with open(file_path, 'r', encoding='utf-8') as file:
                lines = file.readlines()

            # 제목 슬라이드 추가
            title_slide = prs.slides.add_slide(prs.slide_layouts[5])  # 빈 슬라이드
            title_box = title_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            title_frame = title_box.text_frame
            
            title_paragraph = title_frame.add_paragraph()
            title_paragraph.text = filename[:-4]  # 파일 이름 (확장자 제외)
            title_paragraph.font.size = Pt(48)
            title_paragraph.font.name = 'Malgun Gothic'  # 폰트 설정
            title_paragraph.font.color.rgb = RGBColor(255, 255, 0)  # 노란색
            title_paragraph.font.bold = True  # 강조 효과
            title_paragraph.alignment = 1  # 센터 정렬
            title_slide.background.fill.solid()
            title_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검정색

            # 내용 슬라이드 추가
            wrapped_lines = []
            for line in lines:
                wrapped_lines.append(line.strip())

            # 페이지에 추가할 내용
            slides_content = []
            for i in range(0, len(wrapped_lines), 4):
                slides_content.append(wrapped_lines[i:i + 4])  # 4줄씩 묶어서 저장

            # 마지막 페이지 처리
            if len(slides_content) > 1 and len(slides_content[-1]) == 1:
                slides_content[-2].append(slides_content[-1][0])  # 이전 페이지에 추가
                slides_content.pop()  # 마지막 페이지 제거

            # 슬라이드 생성
            for content in slides_content:
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # 빈 슬라이드
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
                text_frame = text_box.text_frame
                
                for j in range(len(content)):
                    p = text_frame.add_paragraph()
                    p.text = content[j][:13]  # 한 줄에 13글자
                    p.font.size = Pt(48)
                    p.font.name = 'Malgun Gothic'  # 폰트 설정
                    p.font.color.rgb = RGBColor(255, 255, 0)  # 노란색
                    p.font.bold = True  # 강조 효과
                    p.alignment = 1  # 센터 정렬

                    # 한 줄 띄우기
                    if j < len(content) - 1:
                        text_frame.add_paragraph()  # 빈 줄 추가

                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검정색

            # 제목 페이지 번호 저장
            titles.append((filename[:-4], len(prs.slides)))  # 파일 제목과 페이지 번호 저장

    # 첫 페이지 생성: 파일 이름과 페이지 번호
    index_slide = prs.slides.add_slide(prs.slide_layouts[5])  # 빈 슬라이드
    index_box = index_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
    index_frame = index_box.text_frame

    index_title = index_frame.add_paragraph()
    index_title.text = "파일 제목 및 페이지 번호"
    index_title.font.size = Pt(48)
    index_title.font.name = 'Malgun Gothic'
    index_title.font.color.rgb = RGBColor(255, 255, 0)  # 노란색
    index_title.font.bold = True
    index_title.alignment = 1  # 센터 정렬

    for title, page_num in titles:
        p = index_frame.add_paragraph()
        p.text = f"{title}: 페이지 {page_num}"
        p.font.size = Pt(32)
        p.font.name = 'Malgun Gothic'
        p.font.color.rgb = RGBColor(255, 255, 0)  # 노란색
        p.font.bold = True
        p.alignment = 1  # 센터 정렬

    index_slide.background.fill.solid()
    index_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검정색

    # PPT 파일 저장
    pptx_file_path = os.path.join(folder_path, 'output.pptx')
    prs.save(pptx_file_path)
    print(f'PPT 파일이 생성되었습니다: {pptx_file_path}')

# 코드 실행
create_ppt_from_text_files()
